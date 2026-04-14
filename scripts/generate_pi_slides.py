"""
Generate the COLM 2026 presentation deck for
"Alignment Correlates of Consensus Susceptibility".

Slide structure mirrors the paper narrative:
  1. Title
  2. Research Question & Motivation
  3. Key Metrics Defined
  4. Experimental Design
  5. Hero Figure — Post-Training Stage Trajectory
  6. BER Heatmap — Detailed Condition × Stage Landscape
  7. SFT Amplification & DPO Repair (takeaways)
  8. Ablation — Pattern Completion vs Social Framing
  9. What the Ablation Tells Us (takeaways)
 10. Cross-Family Conformity Ranking
 11. Behavioral Taxonomy Scatter
 12. Three Behavioral Modes (takeaways)
 13. Scale Bridge — 7B vs 32B
 14. Key Claims (summary)
 15. Limitations & Future Work

Figures are read from:
  - Comparing_Experiments/April_analysis/figures/          (7B within-family)
  - Comparing_Experiments/April_analysis/figures/cross_family/  (cross-family + ablation)
"""
from __future__ import annotations

import os
from dataclasses import dataclass
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parents[1]
FIG_7B = REPO_ROOT / "Comparing_Experiments" / "April_analysis" / "figures"
FIG_CF = FIG_7B / "cross_family"
OUT_DIR = REPO_ROOT / "Comparing_Experiments" / "April_analysis"
OUT_PATH = OUT_DIR / "April_Analysis_Presentation.pptx"

# ── Theme colors ──────────────────────────────────────────────────────
DARK_BLUE = RGBColor(0x1A, 0x3A, 0x5C)
ACCENT_BLUE = RGBColor(0x2E, 0x6D, 0xA4)
BODY_DARK = RGBColor(0x33, 0x33, 0x33)
BODY_LIGHT = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ── Sizes ─────────────────────────────────────────────────────────────
TITLE_SZ = Pt(28)
SUBTITLE_SZ = Pt(16)
SECTION_SZ = Pt(13)
BODY_SZ = Pt(12)
SMALL_SZ = Pt(11)
CAPTION_SZ = Pt(9)


@dataclass(frozen=True)
class Bullet:
    text: str
    level: int = 0


# ── Helpers ───────────────────────────────────────────────────────────

def _set_title(shape, text: str, *, font_size_pt: int = 28) -> None:
    shape.text = text
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if p.runs:
        for r in p.runs:
            r.font.size = Pt(font_size_pt)
            r.font.bold = True
            r.font.color.rgb = DARK_BLUE
    else:
        p.font.size = Pt(font_size_pt)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE


def _set_subtitle(shape, lines: list[str], *, font_size_pt: int = 16) -> None:
    tf = shape.text_frame
    tf.clear()
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(font_size_pt)
        p.font.color.rgb = BODY_LIGHT


def _set_bullets(placeholder, bullets: list[Bullet], *, font_size_pt: int = 12) -> None:
    tf = placeholder.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, b in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = b.text
        p.level = b.level
        p.font.size = Pt(font_size_pt)
        p.font.color.rgb = BODY_DARK


def _add_title_only_slide(prs: Presentation, title: str) -> "pptx.slide.Slide":
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide.shapes.title, title, font_size_pt=24)
    return slide


def _add_image_fit(
    slide,
    image_path: Path,
    *,
    left,
    top,
    width,
    height,
) -> None:
    if not image_path.exists():
        print(f"  WARNING: missing figure {image_path}")
        return
    with Image.open(image_path) as img:
        img_w_px, img_h_px = img.size
    scale = min(width / img_w_px, height / img_h_px)
    w = int(img_w_px * scale)
    h = int(img_h_px * scale)
    x = int(left + (width - w) / 2)
    y = int(top + (height - h) / 2)
    slide.shapes.add_picture(str(image_path), x, y, width=w, height=h)


def _add_textbox(
    slide,
    *,
    left,
    top,
    width,
    height,
    text: str,
    font_size_pt: int = 12,
    bold: bool = False,
    color: RGBColor = BODY_DARK,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.text = text
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.size = Pt(font_size_pt)
    p.font.bold = bold
    p.font.color.rgb = color


def _add_caption(slide, text: str, *, left, top, width) -> None:
    """Small caption text below a figure."""
    _add_textbox(
        slide,
        left=left,
        top=top,
        width=width,
        height=Inches(0.6),
        text=text,
        font_size_pt=9,
        color=BODY_LIGHT,
    )


# ── Slide builders ────────────────────────────────────────────────────

def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_w = prs.slide_width
    slide_h = prs.slide_height
    margin_x = Inches(0.65)
    margin_bottom = Inches(0.45)
    content_top = Inches(1.15)

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 1 — Title
    # ──────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[0])
    _set_title(
        s.shapes.title,
        "Alignment Correlates of Consensus Susceptibility",
        font_size_pt=36,
    )
    subtitle = s.placeholders[1]
    _set_subtitle(
        subtitle,
        [
            "Training-Stage Decomposition and Cross-Family Survey in Large Language Models",
            "",
            "OLMo 7B + 32B (Instruct & Think paths)  |  8 additional model families  |  12 pressure conditions",
            "",
            "Mahdi Saeedi, Jinyi Ye, Luca Luceri  \u2014  University of Southern California",
        ],
        font_size_pt=18,
    )

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 2 — Research Question
    # ──────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[1])
    _set_title(s.shapes.title, "What Are We Studying?", font_size_pt=28)
    body = s.placeholders[1]
    _set_bullets(
        body,
        [
            Bullet("When a model knows the right answer, will it change its mind if fabricated peers all endorse the wrong one?"),
            Bullet("We call this consensus susceptibility \u2014 distinct from sycophancy (agreeing with the user).", level=1),
            Bullet(""),
            Bullet("Why it matters: LLMs are deployed as medical assistants, legal advisors, coding copilots."),
            Bullet("If embedding fake consensus in the prompt flips the answer, that is a concrete deployment vulnerability.", level=1),
            Bullet(""),
            Bullet("Key insight: modern LLMs go through multiple post-training stages (SFT \u2192 DPO \u2192 RL)."),
            Bullet("Do these stages help or hurt? Prior work treated the pipeline as a black box.", level=1),
            Bullet("OLMo is the only major model family releasing every intermediate checkpoint \u2014 so we can look inside.", level=1),
        ],
        font_size_pt=14,
    )

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 3 — Key Metrics
    # ──────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[1])
    _set_title(s.shapes.title, "Key Metrics", font_size_pt=28)
    body = s.placeholders[1]
    _set_bullets(
        body,
        [
            Bullet("BER (Behavioral Error Rate) = wrong-answer endorsement rate"),
            Bullet("The fraction of questions where the model adopts the fabricated wrong answer as its own.", level=1),
            Bullet("Formally: B / N, where B = count of wrong-answer endorsements, N = 400 items per condition.", level=1),
            Bullet(""),
            Bullet("3-State Decomposition: every trial is classified as one of:"),
            Bullet("State A (correct)  |  State B (wrong-answer endorsed)  |  State C (refusal)", level=1),
            Bullet("This matters because a model that refuses everything and one that endorses wrong answers are both 'failing' \u2014 but differently.", level=1),
            Bullet(""),
            Bullet("\u0394 (Delta) = Error_pressure \u2212 Error_control  (the pressure effect)"),
            Bullet("McNemar's exact test with Holm\u2013Bonferroni correction for statistical significance.", level=1),
            Bullet("Wilson 95% confidence intervals on BER.", level=1),
        ],
        font_size_pt=13,
    )

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 4 — Experimental Design
    # ──────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[1])
    _set_title(s.shapes.title, "Experimental Design", font_size_pt=28)
    body = s.placeholders[1]
    _set_bullets(
        body,
        [
            Bullet("Within-family decomposition (OLMo):"),
            Bullet("Instruct path (7B, full coverage): Base \u2192 SFT \u2192 DPO \u2192 Instruct (RLVR), 12 conditions, 6 temperatures", level=1),
            Bullet("Think path (7B): Think-SFT, Think-DPO \u2014 2 temperatures, 4 shared conditions", level=1),
            Bullet("Think path (32B): Think-SFT, Think-DPO, Think \u2014 2 temperatures, 4 shared conditions", level=1),
            Bullet("32B Instruct \u2014 scale comparison on the same shared conditions", level=1),
            Bullet(""),
            Bullet("Cross-family extension:"),
            Bullet("8 additional families: Llama-3-8B, Llama-3.1-70B, Llama-4-Maverick, GPT-4o-Mini, GPT-OSS-20B, Gemini-2.5-Flash-Lite, Grok-4.1-Fast, Claude-Sonnet-4", level=1),
            Bullet("4 core conditions \u00d7 2 temperatures \u00d7 400 items per condition", level=1),
            Bullet(""),
            Bullet("Ablation studies: system-prompt removal + non-social n-gram baseline"),
            Bullet("400 factual QA items across 8 domains (math, science, reasoning, knowledge, etc.)", level=1),
        ],
        font_size_pt=13,
    )

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 5 — Hero Figure: Stage Trajectory
    # ──────────────────────────────────────────────────────────────────
    s = _add_title_only_slide(
        prs,
        "Finding 1: SFT Amplifies Susceptibility, DPO Partially Repairs It",
    )
    _add_image_fit(
        s,
        FIG_7B / "fig_stage_trajectory.png",
        left=margin_x,
        top=content_top,
        width=slide_w - 2 * margin_x,
        height=slide_h - content_top - margin_bottom - Inches(0.6),
    )
    _add_caption(
        s,
        "BER (wrong-answer endorsement rate, B/400) across post-training stages for OLMo-7B. "
        "Instruct path (blue) shows dramatic SFT spike; Think path (orange) remains stable. "
        "Both share the same Base checkpoint.",
        left=margin_x,
        top=slide_h - margin_bottom - Inches(0.5),
        width=slide_w - 2 * margin_x,
    )

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 6 — BER Heatmap
    # ──────────────────────────────────────────────────────────────────
    s = _add_title_only_slide(
        prs,
        "The Full Landscape: BER Across All Stages and Conditions (T=0)",
    )
    _add_image_fit(
        s,
        FIG_7B / "fig_2axis_heatmap_combined.png",
        left=margin_x,
        top=content_top,
        width=slide_w - 2 * margin_x,
        height=slide_h - content_top - margin_bottom - Inches(0.6),
    )
    _add_caption(
        s,
        "Panel A: Instruct-path variants \u00d7 12 conditions, ordered by target-answer repetition count. "
        "Panel B: Think-path variants \u00d7 4 shared conditions, column-aligned. "
        "Color gradient tracks repetition count \u2014 a visual fingerprint of pattern completion.",
        left=margin_x,
        top=slide_h - margin_bottom - Inches(0.5),
        width=slide_w - 2 * margin_x,
    )

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 7 — SFT/DPO Takeaways
    # ──────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[1])
    _set_title(s.shapes.title, "What the Training-Stage Decomposition Reveals", font_size_pt=28)
    body = s.placeholders[1]
    _set_bullets(
        body,
        [
            Bullet("SFT amplifies vulnerability:"),
            Bullet("The SFT checkpoint produces the highest wrong-answer endorsement of any variant \u2014 higher than the raw Base model.", level=1),
            Bullet("SFT simultaneously shows the lowest refusal rate: it suppresses abstention while amplifying pattern following.", level=1),
            Bullet(""),
            Bullet("DPO partially repairs the damage:"),
            Bullet("Endorsement drops and refusal rises relative to SFT after preference optimization.", level=1),
            Bullet("But the final RL stage (Instruct) shows a partial rebound. All variants remain significantly vulnerable.", level=1),
            Bullet(""),
            Bullet("Think path tells the opposite story:"),
            Bullet("SFT reduces endorsement on the Think path. Reasoning training protects against the same pressure.", level=1),
            Bullet("This is invisible in any end-to-end comparison (Base vs final Instruct) \u2014 you need the intermediate checkpoints to see it.", level=1),
        ],
        font_size_pt=13,
    )

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 8 — Ablation Figure
    # ──────────────────────────────────────────────────────────────────
    s = _add_title_only_slide(
        prs,
        "Finding 2: What Looks Like Peer Pressure Is Largely Pattern Completion",
    )
    _add_image_fit(
        s,
        FIG_CF / "fig_ablation_ngram_vs_pressure.png",
        left=Inches(2.0),
        top=content_top,
        width=slide_w - Inches(4.0),
        height=slide_h - content_top - margin_bottom - Inches(0.6),
    )
    _add_caption(
        s,
        "BER under social pressure (blue), without system prompt (orange), and with non-social n-gram baseline (red). "
        "Llama-3.1-70B endorses wrong answers 7.94\u00d7 more from pure pattern repetition than from social framing. "
        "For OLMo-32B, the ratio is 1.12\u00d7 \u2014 social framing adds almost nothing beyond what pattern completion already explains.",
        left=margin_x,
        top=slide_h - margin_bottom - Inches(0.5),
        width=slide_w - 2 * margin_x,
    )

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 9 — Ablation Takeaways
    # ──────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[1])
    _set_title(s.shapes.title, "What the Ablation Tells Us", font_size_pt=28)
    body = s.placeholders[1]
    _set_bullets(
        body,
        [
            Bullet("The Asch-style prompt embeds five repetitions of the wrong answer (\"Participant 1: X ... Participant 5: X\")."),
            Bullet("Is the model responding to social pressure \u2014 or just completing a repetitive pattern?", level=1),
            Bullet(""),
            Bullet("System prompt removal: for some models, endorsement spikes dramatically without the truth-preferring instruction."),
            Bullet("For others, the system prompt provides zero measurable protection.", level=1),
            Bullet(""),
            Bullet("Non-social n-gram baseline: stripping all social framing while keeping the repetitive structure still produces substantial endorsement."),
            Bullet("This reframes the construct: the vulnerability is to repetitive in-context patterns, not to social reasoning.", level=1),
            Bullet(""),
            Bullet("Practical implication: the threat is real regardless of mechanism \u2014 structured repetition in a prompt can flip model answers."),
            Bullet("But the mitigation strategy changes: defend against n-gram repetition, not against social deference.", level=1),
        ],
        font_size_pt=13,
    )

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 10 — Cross-Family Ranking
    # ──────────────────────────────────────────────────────────────────
    s = _add_title_only_slide(
        prs,
        "Finding 3: Susceptibility Varies Massively Across Families",
    )
    _add_image_fit(
        s,
        FIG_CF / "fig_cross_family_headline_ber.png",
        left=margin_x,
        top=content_top,
        width=slide_w - 2 * margin_x,
        height=slide_h - content_top - margin_bottom - Inches(0.6),
    )
    _add_caption(
        s,
        "BER (B/400) on unanimous-confident peer pressure at T=0. "
        "OLMo training-stage checkpoints embedded alongside cross-family models. "
        "Range: 4.5% (Llama-3.1-70B) to 73.8% (OLMo-7B-Instruct-SFT). "
        "Error bars = 95% Wilson CIs.",
        left=margin_x,
        top=slide_h - margin_bottom - Inches(0.5),
        width=slide_w - 2 * margin_x,
    )

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 11 — Behavioral Taxonomy Scatter
    # ──────────────────────────────────────────────────────────────────
    s = _add_title_only_slide(
        prs,
        "Three Distinct Failure Modes Under Identical Pressure",
    )
    _add_image_fit(
        s,
        FIG_CF / "fig4_refusal_endorsement.png",
        left=Inches(1.5),
        top=content_top,
        width=slide_w - Inches(3.0),
        height=slide_h - content_top - margin_bottom - Inches(0.6),
    )
    _add_caption(
        s,
        "\u0394 endorsement (BER change) vs \u0394 refusal rate at T=0. "
        "Endorsement-dominant (right): model adopts wrong answers. "
        "Refusal-dominant (top): model declines to answer. "
        "Context-insensitive (origin): model ignores the pressure entirely.",
        left=margin_x,
        top=slide_h - margin_bottom - Inches(0.5),
        width=slide_w - 2 * margin_x,
    )

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 12 — Three Behavioral Modes (text)
    # ──────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[1])
    _set_title(s.shapes.title, "The Three Behavioral Modes", font_size_pt=28)
    body = s.placeholders[1]
    _set_bullets(
        body,
        [
            Bullet("Endorsement-dominant (GPT-4o-Mini, OLMo-32B-Instruct):"),
            Bullet("Model adopts the fabricated consensus as its own answer at high rates while rarely refusing.", level=1),
            Bullet("This is the most dangerous failure mode \u2014 the model confidently gives users wrong answers.", level=1),
            Bullet(""),
            Bullet("Refusal-dominant (Llama-3-8B):"),
            Bullet("Model triggers safety-aligned refusal rather than endorsing. Still failing, but differently.", level=1),
            Bullet("Under fixed-N design, refusals inflate error rates but do not represent wrong-answer adoption.", level=1),
            Bullet(""),
            Bullet("Context-insensitive (Claude-Sonnet-4, Grok-4.1-Fast, Llama-3.1-70B):"),
            Bullet("Model answers as if the fabricated peers were not there. Maintains baseline accuracy.", level=1),
            Bullet("All reasoning-capable models in our survey show non-significant peer conformity.", level=1),
            Bullet(""),
            Bullet("Alignment approach, not just capability, determines which mode a model falls into."),
        ],
        font_size_pt=13,
    )

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 13 — Scale Bridge
    # ──────────────────────────────────────────────────────────────────
    s = _add_title_only_slide(
        prs,
        "Scale Comparison: OLMo 7B vs 32B",
    )
    _add_image_fit(
        s,
        FIG_CF / "fig_scale_bridge.png",
        left=Inches(1.5),
        top=content_top,
        width=slide_w - Inches(3.0),
        height=slide_h - content_top - margin_bottom - Inches(0.6),
    )
    _add_caption(
        s,
        "BER at each training stage for OLMo-7B vs 32B on the 4 shared conditions at T=0. "
        "The Think path at 32B shows lower susceptibility than Instruct, consistent with the 7B finding.",
        left=margin_x,
        top=slide_h - margin_bottom - Inches(0.5),
        width=slide_w - 2 * margin_x,
    )

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 14 — Key Claims Summary
    # ──────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[1])
    _set_title(s.shapes.title, "What We Can Safely Claim", font_size_pt=28)
    body = s.placeholders[1]
    _set_bullets(
        body,
        [
            Bullet("1. Training-stage decomposition (novel):"),
            Bullet("SFT amplifies wrong-answer endorsement; DPO partially reverses it. These opposing effects cancel in any end-to-end comparison.", level=1),
            Bullet("The Think (reasoning) path is protected from the same pressure that devastates the Instruct path.", level=1),
            Bullet(""),
            Bullet("2. Pattern completion, not social deference (novel):"),
            Bullet("Stripping social framing while preserving repetitive structure still produces substantial endorsement.", level=1),
            Bullet("The vulnerability is to autoregressive pattern completion, not to sociological reasoning.", level=1),
            Bullet(""),
            Bullet("3. Three behavioral modes, not a spectrum (novel):"),
            Bullet("Models fail under pressure in categorically different ways: endorse, refuse, or resist.", level=1),
            Bullet("Alignment approach, not capability alone, shapes which mode a model falls into.", level=1),
        ],
        font_size_pt=14,
    )

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 15 — Limitations & Future Work
    # ──────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[1])
    _set_title(s.shapes.title, "Limitations & Future Work", font_size_pt=28)
    body = s.placeholders[1]
    _set_bullets(
        body,
        [
            Bullet("Limitations:"),
            Bullet("Cannot prove causation for closed models whose training details are proprietary.", level=1),
            Bullet("The effect is tied to the structured n-gram format; generalization to multi-turn conversational pressure is untested.", level=1),
            Bullet("DPO is applied to SFT weights \u2014 the observed mitigation reflects a sequential trajectory, not an independent counterfactual.", level=1),
            Bullet("Reasoning models may resist via structural interrupt (<think> token) rather than deliberation.", level=1),
            Bullet(""),
            Bullet("Future work:"),
            Bullet("Multi-turn pressure: does sustained conversational consensus produce the same effect?", level=1),
            Bullet("Causal probing: intervene on repetition-sensitive attention heads to test the pattern-completion hypothesis.", level=1),
            Bullet("Mitigation design: n-gram diversity requirements in context windows.", level=1),
        ],
        font_size_pt=13,
    )

    return prs


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = build_deck()
    prs.save(str(OUT_PATH))
    print(f"Wrote: {OUT_PATH.relative_to(REPO_ROOT)}")


if __name__ == "__main__":
    main()
