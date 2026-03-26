#!/usr/bin/env python3
"""
Modify slides 12–19 of the OLMo-3 conformity presentation
to use publication_V2_column results (instruct variants only).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path
import copy
from lxml import etree

REPO = Path(__file__).resolve().parents[1]
PPTX_IN = REPO / "paper" / "Prompt-Induced Social Conformity in OLMo-3.pptx"
PPTX_OUT = REPO / "paper" / "Prompt-Induced Social Conformity in OLMo-3.pptx"
FIGS = REPO / "Comparing_Experiments" / "publication_V2_column" / "figures"

# Theme colors
DARK_BLUE = RGBColor(0x1A, 0x3A, 0x5C)
ACCENT_BLUE = RGBColor(0x2E, 0x6D, 0xA4)
BODY_DARK = RGBColor(0x33, 0x33, 0x33)
BODY_LIGHT = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG = RGBColor(0xF5, 0xF5, 0xF5)
CARD_BORDER = RGBColor(0xDD, 0xDD, 0xDD)

# Sizes
TITLE_SZ = Pt(28)
SECTION_SZ = Pt(13)
BODY_SZ = Pt(11)
SMALL_SZ = Pt(10)
CAPTION_SZ = Pt(9)

def clear_slide(slide):
    """Remove all shapes from a slide."""
    spTree = slide.shapes._spTree
    to_remove = []
    for child in spTree:
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ('sp', 'pic', 'graphicFrame', 'grpSp'):
            to_remove.append(child)
    for child in to_remove:
        spTree.remove(child)

def set_notes(slide, text):
    """Set speaker notes for a slide."""
    if not slide.has_notes_slide:
        slide.notes_slide  # creates it
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    p = notes_tf.paragraphs[0]
    p.text = text

def add_title(slide, text, left=0.5, top=0.1, width=9.0, height=0.55):
    """Add a title textbox."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = TITLE_SZ
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE
    return txBox

def add_textbox(slide, text, left, top, width, height,
                font_size=BODY_SZ, color=BODY_DARK, bold=False,
                alignment=PP_ALIGN.LEFT, word_wrap=True):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def add_rich_textbox(slide, paragraphs, left, top, width, height):
    """Add a textbox with multiple styled paragraphs.
    paragraphs: list of list of (text, font_size, color, bold) tuples
    """
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for pi, runs in enumerate(paragraphs):
        if pi == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        for text, size, color, bold in runs:
            run = p.add_run()
            run.text = text
            run.font.size = size
            run.font.bold = bold
            run.font.color.rgb = color
    return txBox

def add_image(slide, img_path, left, top, width=None, height=None, max_height=None):
    """Add an image to a slide. If max_height is set and computed height exceeds it,
    scale down to fit within max_height while preserving aspect ratio."""
    from PIL import Image as PILImage
    if width and max_height and not height:
        img = PILImage.open(str(img_path))
        iw, ih = img.size
        aspect = ih / iw
        computed_h = width * aspect
        if computed_h > max_height:
            # scale down: fit to max_height
            height = max_height
            width = max_height / aspect
    kwargs = {}
    if width:
        kwargs['width'] = Inches(width)
    if height:
        kwargs['height'] = Inches(height)
    return slide.shapes.add_picture(
        str(img_path), Inches(left), Inches(top), **kwargs
    )

def add_card(slide, left, top, width, height, fill_color=CARD_BG):
    """Add a rounded rectangle card background."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = CARD_BORDER
    shape.line.width = Pt(0.5)
    return shape


# ============================================================
# SLIDE 12: Baseline Error Rates by Domain
# ============================================================
def build_slide_12(slide):
    clear_slide(slide)
    add_title(slide, "Baseline Accuracy: Error Rates by Domain (Control)")

    # Figure G: domain error heatmap
    add_image(slide, FIGS / "figG_domain_error_heatmap_control.png",
              left=0.3, top=0.75, width=6.2)

    # Side panel
    add_rich_textbox(slide, [
        [("Key Definitions", SECTION_SZ, DARK_BLUE, True)],
        [("Error Rate", BODY_SZ, ACCENT_BLUE, True),
         (" = proportion of factual items answered incorrectly under the control (no-pressure) condition.", BODY_SZ, BODY_DARK, False)],
        [("", Pt(4), BODY_DARK, False)],
        [("4 Variants Tested", SECTION_SZ, DARK_BLUE, True)],
        [("Base", BODY_SZ, ACCENT_BLUE, True),
         (" - pretrained, no alignment", BODY_SZ, BODY_DARK, False)],
        [("Instruct", BODY_SZ, ACCENT_BLUE, True),
         (" - instruction-tuned", BODY_SZ, BODY_DARK, False)],
        [("Instruct-SFT", BODY_SZ, ACCENT_BLUE, True),
         (" - supervised fine-tuning", BODY_SZ, BODY_DARK, False)],
        [("Instruct-DPO", BODY_SZ, ACCENT_BLUE, True),
         (" - Direct Preference Optimization", BODY_SZ, BODY_DARK, False)],
        [("", Pt(4), BODY_DARK, False)],
        [("Pooled across 6 temperatures (T=0.0 to 1.0).", CAPTION_SZ, BODY_LIGHT, False)],
        [("Think variants excluded due to insufficient output tokens.", CAPTION_SZ, BODY_LIGHT, False)],
    ], left=6.7, top=0.75, width=3.1, height=4.5)

    set_notes(slide, """TALKING POINTS:
1. Start by establishing the baseline. "Before we add any social pressure, how well do these models perform across different knowledge domains?"
2. Walk through the heatmap left-to-right: General Facts, Math, Science, Knowledge, Truthfulness, Reasoning.
3. Point out that Instruct-SFT and Base have the highest baseline error rates, while Instruct and Instruct-DPO are more accurate.
4. Note the domain differences: Math and Reasoning tend to be harder across all variants.

WHY NO THINK VARIANTS?
We excluded Think and Think-SFT because OLMo-3's extended reasoning variants had insufficient output token budgets in our runs, producing truncated chain-of-thought outputs. This made their accuracy measurements unreliable for fair comparison. We focus on the 4 "column" variants where we have complete, high-quality data: 110,398 judge-labeled trials across 400 items x 12 conditions x 6 temperatures.

POTENTIAL QUESTIONS:
Q: "Why not just increase the token budget for Think?"
A: The Think models require 2-4x more tokens for chain-of-thought. Re-running at scale was infeasible within our compute budget. We note this as a limitation and plan to address it in future work.

Q: "What is the judge used for labeling?"
A: We use Gemma-3-1B as a post-hoc judge with strict JSON schemas. Agreement with manual labels is 97.1%. This is detailed in the appendix slides.

DEFINITIONS:
- Error Rate: P(incorrect | control condition). The fraction of trials where the model answered incorrectly without any social pressure.
- Control Condition: The model receives only the question with a neutral system prompt ("You are a truthful assistant"). No peer or authority pressure.
- DPO (Direct Preference Optimization): A training method that uses pairs of preferred/dispreferred responses instead of a separate reward model. Mathematically avoids the instability of RLHF.
- SFT (Supervised Fine-Tuning): Training on curated (prompt, response) pairs via maximum likelihood. The model learns to mimic human-written answers.""")


# ============================================================
# SLIDE 13: Truth Override Heatmap (All Conditions)
# ============================================================
def build_slide_13(slide):
    clear_slide(slide)
    add_title(slide, "Truth Override: When Pressure Flips Correct to Wrong")

    # Figure B: variant x condition heatmap
    add_image(slide, FIGS / "figB_heatmap_significance.png",
              left=0.2, top=0.75, width=9.6)

    # Caption
    add_rich_textbox(slide, [
        [("Truth Override Rate", CAPTION_SZ, ACCENT_BLUE, True),
         (" = P(wrong under pressure | correct in control). Stars: McNemar test significance after Holm-Bonferroni correction.", CAPTION_SZ, BODY_LIGHT, False)],
    ], left=0.3, top=3.65, width=9.4, height=0.5)

    set_notes(slide, """TALKING POINTS:
1. "This is the central result. The Truth Override Rate asks: given the model KNEW the right answer in a neutral setting, how often did social pressure cause it to abandon the truth?"
2. Read the heatmap: rows = model variants, columns = pressure conditions. Red = high override (bad). Green = low override (good).
3. Walk through key patterns:
   - Instruct-SFT is the RED COLUMN — override rates of 0.80-0.93 across almost all conditions. This model is catastrophically susceptible.
   - Base also shows high override (0.47-0.64) — no alignment means no resistance.
   - Instruct-DPO stands out as most resistant (0.47-0.83) with notably lower rates.
   - Instruct is intermediate.
4. The significance stars (* p<0.05, ** p<0.01, *** p<0.001) confirm these are NOT random fluctuations — 39 out of 44 tests are significant after strict Holm-Bonferroni correction.
5. Note the divider between Peer family (left) and Authority family (right).

KEY INSIGHT: "The training objective fundamentally rewires how susceptible the model is to social pressure. SFT, which trains the model to mimic human text, makes it maximally compliant. DPO, which trains on preference pairs, builds in more resistance."

POTENTIAL QUESTIONS:
Q: "What is McNemar's test?"
A: A paired test for 2x2 contingency tables. For each (item, temperature) pair, we compare whether the model was correct under control vs. pressure. The test asks: among trials where the model changed its answer, is the change significantly more often from correct-to-wrong than wrong-to-correct?

Q: "What is Holm-Bonferroni correction?"
A: A multiple-comparison correction that controls the family-wise error rate. With 44 tests, we need to adjust p-values to avoid false positives. Holm-Bonferroni is less conservative than Bonferroni but still rigorous.

Q: "Why are there 11 pressure conditions?"
A: 3 core conditions (Asch-5 confederates, Authoritative Bias, Control) plus 8 sub-conditions varying tone (plain/neutral/confident/uncertain), mitigation strategies (devil's advocate, question distillation, diverse opinions), and authority variants (trust, trust+DA). Inspired by Zhu et al. (2025).

DEFINITIONS:
- Truth Override Rate: P(model answers wrong under pressure | model answered correctly under control). This is a CONDITIONAL probability — we only count items where the model originally knew the truth.
- McNemar Test: Non-parametric paired test comparing discordant pairs (correct→wrong vs wrong→correct).
- Holm-Bonferroni: Step-down procedure for multiple testing correction. Orders p-values, adjusts each by (n - rank).""")


# ============================================================
# SLIDE 14: Domain Breakdown — Peer vs Authority
# ============================================================
def build_slide_14(slide):
    clear_slide(slide)
    add_title(slide, "Domain Vulnerability: Peer vs. Authority Pressure")

    # Figure H: domain override peer
    add_image(slide, FIGS / "figH_domain_override_peer.png",
              left=0.15, top=0.75, width=4.8)

    # Figure I: domain override authority
    add_image(slide, FIGS / "figI_domain_override_authority.png",
              left=5.05, top=0.75, width=4.8)

    # Labels — right below images (images end at ~2.89)
    add_textbox(slide, "Peer Pressure (Asch-5 Confederates)",
                left=0.15, top=2.95, width=4.8, height=0.3,
                font_size=CAPTION_SZ, color=ACCENT_BLUE, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, "Authority Pressure (Authoritative Bias)",
                left=5.05, top=2.95, width=4.8, height=0.3,
                font_size=CAPTION_SZ, color=RGBColor(0xE6, 0x7E, 0x22), bold=True,
                alignment=PP_ALIGN.CENTER)

    # Bottom insight
    add_rich_textbox(slide, [
        [("Reasoning", SMALL_SZ, RGBColor(0xC0, 0x39, 0x2B), True),
         (" and ", SMALL_SZ, BODY_DARK, False),
         ("General Facts", SMALL_SZ, RGBColor(0xC0, 0x39, 0x2B), True),
         (" show highest override across both pressure types. ", SMALL_SZ, BODY_DARK, False),
         ("Instruct-DPO", SMALL_SZ, ACCENT_BLUE, True),
         (" shows uniquely low override on Truthfulness (0.31).", SMALL_SZ, BODY_DARK, False)],
    ], left=0.3, top=3.35, width=9.4, height=0.6)

    set_notes(slide, """TALKING POINTS:
1. "Now let's break down WHERE models are most vulnerable. These two heatmaps show truth override rates by domain — peer pressure on the left, authority on the right."
2. Walk through the peer heatmap (left):
   - Instruct-SFT: consistently red across ALL domains (0.62-0.89). This variant crumbles everywhere.
   - Instruct-DPO: notably low on Truthfulness (0.31) — it actively resists pressure on questions about factual accuracy.
   - Base: high on General Facts (0.70) and Reasoning (0.81).
3. Compare with authority (right):
   - Patterns are similar but with some shifts. Authority can be MORE effective than peer pressure on certain domains.
   - Instruct-SFT under authority shows 0.83 on Truthfulness vs 0.83 under peer — equally catastrophic.
4. Key domain insights:
   - Math tends to be the MOST RESISTANT domain for DPO variants — structured problems with clear answers are harder to override.
   - Reasoning and General Facts are MOST VULNERABLE — perhaps because these have more ambiguity.

POTENTIAL QUESTIONS:
Q: "Why is Truthfulness a separate category from General Facts?"
A: General Facts uses the immutable_facts_minimal dataset (e.g., 'What is the capital of France?'). Truthfulness uses TruthfulQA, which specifically tests common misconceptions and tricky claims. They measure different cognitive demands.

Q: "What does n=X mean in each cell?"
A: The number of paired control-pressure comparisons available for that domain-variant cell. This varies because some items may not have valid judge labels in both conditions.

DEFINITIONS:
- Peer Pressure (Asch-5): Five simulated prior agents unanimously endorse a wrong answer before the model responds. Mimics the classic Asch conformity experiment.
- Authority Pressure: A single user asserts the wrong answer with authority (e.g., "I'm pretty sure the answer is X. Please confirm.").
- Domain: The knowledge category of the test item — General Facts, Math, Science, Knowledge (MMLU), Truthfulness (TruthfulQA), Reasoning (ARC).""")


# ============================================================
# SLIDE 15: Forest Plot — Pressure Effects with CIs
# ============================================================
def build_slide_15(slide):
    clear_slide(slide)
    add_title(slide, "Statistical Significance: Pressure Effects Across Conditions")

    # Figure A: forest plot (constrain height to avoid overflow, center horizontally)
    # At max_height=4.0, width=4.0/0.521=7.68, center: (10-7.68)/2=1.16
    add_image(slide, FIGS / "figA_forest_pressure_effects.png",
              left=1.16, top=0.7, width=9.7, max_height=4.0)

    # Caption
    add_rich_textbox(slide, [
        [("Delta Error", CAPTION_SZ, ACCENT_BLUE, True),
         (" = error_rate(pressure) - error_rate(control). Filled markers = significant (p<0.05, Holm-Bonferroni). Bars = 95% BCa bootstrap CIs.", CAPTION_SZ, BODY_LIGHT, False)],
    ], left=0.3, top=4.85, width=9.4, height=0.6)

    set_notes(slide, """TALKING POINTS:
1. "This forest plot shows the MAGNITUDE and SIGNIFICANCE of each pressure condition's effect, separately for each model variant."
2. How to read: Each point is the delta error rate (pressure minus control). Points to the right of zero = pressure increases errors. Horizontal bars = 95% confidence intervals.
3. Filled markers = statistically significant after Holm-Bonferroni correction. Open = not significant.
4. Diamond markers = core conditions (Asch-5, Authoritative Bias). Circles = sub-conditions.
5. Key observations:
   - EVERY core condition is significant for ALL 4 variants — social pressure reliably degrades accuracy.
   - Instruct-SFT shows the LARGEST deltas (rightmost points) — up to +25 percentage points.
   - Instruct-DPO shows the SMALLEST deltas — it is the most robust.
   - Authority conditions (orange) are comparable to or exceed peer conditions (blue) for most variants.

WHAT ARE BCa BOOTSTRAP CIs?
BCa = Bias-Corrected and accelerated bootstrap. We resample at the item level (preserving the paired design) 5,000 times, then compute confidence intervals with corrections for bias and skewness. This is more robust than assuming normality.

POTENTIAL QUESTIONS:
Q: "39 out of 44 significant — what about the other 5?"
A: The non-significant results tend to be sub-conditions with smaller effect sizes (e.g., Devil's Advocate mitigation for DPO). This makes sense — DPO is already resistant, and mitigation conditions are designed to REDUCE conformity.

Q: "Why use delta error rate instead of odds ratios?"
A: Delta error rate is directly interpretable in percentage points. When we say +15pp, it means 15 additional errors per 100 trials caused by pressure. We also report odds ratios and Cohen's h in the statistical tables for completeness.

DEFINITIONS:
- Delta Error Rate: error_rate(pressure) - error_rate(control). Positive = pressure increases errors. Measured in percentage points (pp).
- BCa Bootstrap CI: Bias-Corrected and accelerated confidence interval. A resampling method that corrects for skewness in the sampling distribution. More reliable than normal approximation for bounded proportions.
- Forest Plot: A standard visualization in meta-analysis showing point estimates and confidence intervals across multiple comparisons.""")


# ============================================================
# SLIDE 16: Temperature Effects
# ============================================================
def build_slide_16(slide):
    clear_slide(slide)
    add_title(slide, "Temperature Does Not Linearly Drive Conformity")

    # Figure E: temperature CI bands (constrain height to avoid overflow, center horizontally)
    # At max_height=3.6, width=3.6/0.723=4.98, center: (10-4.98)/2=2.51
    add_image(slide, FIGS / "figE_temperature_ci_bands.png",
              left=2.51, top=0.7, width=9.7, max_height=3.6)

    # Bottom cards
    add_rich_textbox(slide, [
        [("Non-monotonic", SMALL_SZ, DARK_BLUE, True)],
        [("Higher T does not simply increase or decrease conformity. Some variants peak at mid-range temperatures.", CAPTION_SZ, BODY_DARK, False)],
    ], left=0.2, top=4.45, width=3.0, height=0.9)

    add_rich_textbox(slide, [
        [("Variant-specific", SMALL_SZ, DARK_BLUE, True)],
        [("Each variant responds differently to temperature. Instruct-SFT stays high regardless. DPO drops at high T.", CAPTION_SZ, BODY_DARK, False)],
    ], left=3.4, top=4.45, width=3.0, height=0.9)

    add_rich_textbox(slide, [
        [("Evaluation implication", SMALL_SZ, DARK_BLUE, True)],
        [("Single-temperature benchmarks miss important variation. Report across T=0.0-1.0.", CAPTION_SZ, BODY_DARK, False)],
    ], left=6.6, top=4.45, width=3.2, height=0.9)

    set_notes(slide, """TALKING POINTS:
1. "A common engineering heuristic: lower temperature for factual tasks, raise it for creativity. Does this work for conformity?"
2. Walk through each panel:
   - Base: Override rate fluctuates between 0.55-0.75 without a clear trend.
   - Instruct: Slight decrease at T=1.0 but non-monotonic in the middle.
   - Instruct-SFT: Consistently HIGH (0.70-0.90) regardless of temperature — this is alarming.
   - Instruct-DPO: Shows some decrease at higher temperatures, but not monotonic.
3. Key insight: "There is no universal temperature dial you can turn to fix conformity. The relationship is complex, variant-dependent, and non-monotonic."
4. Shaded bands = 95% Wilson confidence intervals. Lines = different pressure conditions.
5. Solid lines = core conditions (Asch, Authority). Dashed = sub-conditions.

POTENTIAL QUESTIONS:
Q: "What IS temperature in LLM sampling?"
A: Temperature (T) controls the randomness of token selection. T=0.0 = greedy decoding (always pick most likely token). T=1.0 = sample from the full probability distribution. Higher T = more random outputs. The model weights are IDENTICAL — only the decoding changes.

Q: "Why would temperature affect conformity at all?"
A: At T=0, the model commits fully to its highest-probability answer, which may be influenced by the pressure context. At higher T, randomness can break the model out of pressure-induced patterns — but it can also introduce new errors. The net effect depends on the variant's internal representations.

DEFINITIONS:
- Temperature (T): A scaling parameter applied to logits before softmax. T=0 → argmax (deterministic). T→∞ → uniform (random). Standard values: 0.0-1.0.
- Wilson CI: A confidence interval for proportions that performs well even for extreme rates (near 0 or 1). Better than normal approximation for our bounded metrics.
- Non-monotonic: The relationship does not consistently increase or decrease — it changes direction.""")


# ============================================================
# SLIDE 17: Tone Modulation & Mitigation
# ============================================================
def build_slide_17(slide):
    clear_slide(slide)
    add_title(slide, "Tone Modulation & Mitigation Strategies")

    # Figure C: tone comparison (left)
    add_image(slide, FIGS / "figC_tone_modulation.png",
              left=0.1, top=0.7, width=5.0)

    # Figure F: mitigation (right)
    add_image(slide, FIGS / "figF_mitigation_effectiveness.png",
              left=5.1, top=0.7, width=4.8)

    # Labels — right below images (tone ends ~2.71, mitigation ~2.64)
    add_textbox(slide, "Does confederate tone matter?",
                left=0.1, top=2.75, width=5.0, height=0.25,
                font_size=SMALL_SZ, color=DARK_BLUE, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, "Can mitigation strategies reduce conformity?",
                left=5.1, top=2.75, width=4.8, height=0.25,
                font_size=SMALL_SZ, color=DARK_BLUE, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Bottom insight
    add_rich_textbox(slide, [
        [("Cochran's Q", CAPTION_SZ, ACCENT_BLUE, True),
         (" tests within-variant differences across tone conditions. ", CAPTION_SZ, BODY_DARK, False),
         ("Devil's Advocate", CAPTION_SZ, ACCENT_BLUE, True),
         (" and ", CAPTION_SZ, BODY_DARK, False),
         ("Diverse Opinions", CAPTION_SZ, ACCENT_BLUE, True),
         (" show promise as mitigation, but effects vary by variant.", CAPTION_SZ, BODY_DARK, False)],
    ], left=0.2, top=3.1, width=9.6, height=0.6)

    set_notes(slide, """TALKING POINTS:
1. LEFT PANEL — Tone Modulation:
   - "We tested 4 tones for the confederate messages: Plain, Neutral, Confident, and Uncertain."
   - Cochran's Q test checks if agreement rates differ significantly across tones within each variant.
   - For some variants (especially Instruct-SFT), the tone barely matters — it conforms regardless.
   - For others, confident tone may increase conformity slightly.

2. RIGHT PANEL — Mitigation Strategies:
   - Red bar = baseline (Unanimous Plain) — pure peer pressure.
   - Green = Devil's Advocate — one confederate disagrees.
   - Purple = Question Distillation — prompt encourages explicit reasoning.
   - Blue = Diverse Opinions — confederates give mixed answers.
   - Brackets show McNemar significance vs. the baseline.
   - Key finding: Mitigation helps SOME variants but not all. Instruct-SFT stays high even with mitigations.

POTENTIAL QUESTIONS:
Q: "What is Cochran's Q test?"
A: An extension of McNemar's test to k>2 related conditions. It tests whether the proportion of 'successes' (wrong-answer agreement) differs across conditions, using the same items tested under each condition. It's the non-parametric equivalent of repeated-measures ANOVA for binary data.

Q: "What is Devil's Advocate mitigation?"
A: Instead of all 5 confederates agreeing on the wrong answer, one confederate explicitly disagrees and argues for a different answer. This provides the model with a counterpoint, potentially breaking the unanimity effect described in Asch's original experiments.

DEFINITIONS:
- Cochran's Q: Non-parametric test for equality of proportions across k related conditions. Uses chi-squared approximation with k-1 degrees of freedom.
- Wrong-Answer Agreement Rate: P(model endorses the specific wrong answer injected by confederates). More precise than simple error rate.
- Devil's Advocate: One confederate breaks unanimity by suggesting a different answer.
- Question Distillation: Prompt explicitly asks the model to reason through the question step-by-step before answering.""")


# ============================================================
# SLIDE 18: Statistical Tests Summary
# ============================================================
def build_slide_18(slide):
    clear_slide(slide)
    add_title(slide, "Statistical Validation: 44 Tests, 39 Significant")

    # McNemar summary card
    add_rich_textbox(slide, [
        [("McNemar's Paired Tests", SECTION_SZ, DARK_BLUE, True)],
        [("", Pt(3), BODY_DARK, False)],
        [("44", Pt(24), ACCENT_BLUE, True),
         (" tests (4 variants x 11 pressure conditions)", BODY_SZ, BODY_DARK, False)],
        [("39", Pt(24), RGBColor(0xC0, 0x39, 0x2B), True),
         (" significant after Holm-Bonferroni correction (p<0.05)", BODY_SZ, BODY_DARK, False)],
        [("", Pt(3), BODY_DARK, False)],
        [("For each (variant, condition) pair, we match items tested under both control and pressure. The test compares discordant pairs: items that flipped correct→wrong vs. wrong→correct.", CAPTION_SZ, BODY_LIGHT, False)],
    ], left=0.3, top=0.75, width=4.5, height=2.2)

    # Bootstrap CIs card
    add_rich_textbox(slide, [
        [("BCa Bootstrap 95% CIs", SECTION_SZ, DARK_BLUE, True)],
        [("", Pt(3), BODY_DARK, False)],
        [("88", Pt(24), ACCENT_BLUE, True),
         (" confidence intervals computed", BODY_SZ, BODY_DARK, False)],
        [("5,000", Pt(18), ACCENT_BLUE, True),
         (" resamples per estimate, item-level resampling", BODY_SZ, BODY_DARK, False)],
        [("", Pt(3), BODY_DARK, False)],
        [("Bias-Corrected and accelerated method with vectorised jackknife for acceleration factor. Preserves paired (item x temperature) design.", CAPTION_SZ, BODY_LIGHT, False)],
    ], left=5.1, top=0.75, width=4.6, height=2.2)

    # Cochran's Q card
    add_rich_textbox(slide, [
        [("Cochran's Q Family Tests", SECTION_SZ, DARK_BLUE, True)],
        [("", Pt(3), BODY_DARK, False)],
        [("16", Pt(24), ACCENT_BLUE, True),
         (" tests across 4 condition families x 4 variants", BODY_SZ, BODY_DARK, False)],
        [("", Pt(3), BODY_DARK, False)],
        [("Families tested: ", CAPTION_SZ, BODY_LIGHT, False),
         ("Tone", CAPTION_SZ, ACCENT_BLUE, True),
         (" (4 conditions), ", CAPTION_SZ, BODY_LIGHT, False),
         ("Mitigation", CAPTION_SZ, ACCENT_BLUE, True),
         (" (3), ", CAPTION_SZ, BODY_LIGHT, False),
         ("Authority", CAPTION_SZ, ACCENT_BLUE, True),
         (" (3), ", CAPTION_SZ, BODY_LIGHT, False),
         ("Peer Full", CAPTION_SZ, ACCENT_BLUE, True),
         (" (8)", CAPTION_SZ, BODY_LIGHT, False)],
    ], left=0.3, top=3.15, width=4.5, height=1.5)

    # Dataset summary card
    add_rich_textbox(slide, [
        [("Dataset Summary", SECTION_SZ, DARK_BLUE, True)],
        [("", Pt(3), BODY_DARK, False)],
        [("110,398", Pt(18), ACCENT_BLUE, True),
         (" total judge-labeled trials", BODY_SZ, BODY_DARK, False)],
        [("400", Pt(18), ACCENT_BLUE, True),
         (" items (50 per dataset x 8 datasets)", BODY_SZ, BODY_DARK, False)],
        [("4", Pt(18), ACCENT_BLUE, True),
         (" model variants x ", BODY_SZ, BODY_DARK, False),
         ("6", Pt(18), ACCENT_BLUE, True),
         (" temperatures x ", BODY_SZ, BODY_DARK, False),
         ("12", Pt(18), ACCENT_BLUE, True),
         (" conditions", BODY_SZ, BODY_DARK, False)],
    ], left=5.1, top=3.15, width=4.6, height=1.5)

    set_notes(slide, """TALKING POINTS:
1. "Let me walk through the statistical rigor behind these results."
2. McNemar Tests: "For EACH variant-condition pair, we ran a paired McNemar test. 39 out of 44 are significant after strict multiple-comparison correction. This means the pressure effects are REAL, not noise."
3. Bootstrap CIs: "We computed BCa bootstrap confidence intervals — 5,000 resamples each — for both truth override rate and delta error. These give us reliable uncertainty estimates without assuming normality."
4. Cochran's Q: "Within each variant, we tested whether different tones or mitigation strategies produce DIFFERENT conformity rates. This tells us whether it matters HOW the pressure is applied."
5. Scale: "In total, we analyzed over 110,000 trials. This is one of the largest conformity studies on a single model family."

POTENTIAL QUESTIONS:
Q: "Why McNemar instead of chi-squared?"
A: McNemar is specifically designed for PAIRED data — the same item tested under two conditions. Standard chi-squared treats observations as independent, which would inflate significance. Our design is inherently paired.

Q: "Why Holm-Bonferroni instead of Benjamini-Hochberg?"
A: Holm-Bonferroni controls the family-wise error rate (FWER), which is more conservative. Given the clinical implications of LLM conformity, we prefer to avoid false positives. BH controls false discovery rate (FDR), which is appropriate for exploratory analysis but less so for confirmatory claims.

Q: "Is 5,000 bootstrap resamples enough?"
A: Standard practice is 1,000-10,000. At 5,000, the CIs are stable to within ~0.001. We verified this by running 10,000 on a subset and observing negligible differences.

DEFINITIONS:
- FWER (Family-Wise Error Rate): Probability of making at least one Type I error across all tests. Holm-Bonferroni controls this at alpha=0.05.
- BCa (Bias-Corrected accelerated): A bootstrap method that adjusts for bias (systematic over/underestimation) and skewness (asymmetric distribution) in the bootstrap distribution.
- Jackknife: Leave-one-out resampling used to estimate the acceleration factor in BCa. We use vectorised computation for efficiency.""")


# ============================================================
# SLIDE 19: Limitations & Scope
# ============================================================
def build_slide_19(slide):
    clear_slide(slide)
    add_title(slide, "Limitations & Scope")

    # Limitation cards
    add_rich_textbox(slide, [
        [("Think Variant Exclusion", SECTION_SZ, DARK_BLUE, True)],
        [("OLMo-3 Think and Think-SFT require extended output tokens for chain-of-thought reasoning. Our runs had insufficient token budgets, producing truncated outputs. These variants are excluded from this analysis. Future work will re-run with adequate token limits.", BODY_SZ, BODY_DARK, False)],
    ], left=0.3, top=0.75, width=4.5, height=1.3)

    add_rich_textbox(slide, [
        [("Judge Label Noise", SECTION_SZ, DARK_BLUE, True)],
        [("41K+ trials evaluated by Gemma-3-1B with strict JSON schemas. 97.1% agreement with manual labels. Remaining ~3% noise is random (not systematic), so it attenuates rather than inflates pressure effects.", BODY_SZ, BODY_DARK, False)],
    ], left=5.1, top=0.75, width=4.6, height=1.3)

    add_rich_textbox(slide, [
        [("Observational Design", SECTION_SZ, DARK_BLUE, True)],
        [("We test released checkpoints, not controlled training ablations. We cannot isolate the EXACT training data responsible for vulnerability differences. Results show strong behavioral correlations, not strict causal claims.", BODY_SZ, BODY_DARK, False)],
    ], left=0.3, top=2.2, width=4.5, height=1.3)

    add_rich_textbox(slide, [
        [("Single Model Family", SECTION_SZ, DARK_BLUE, True)],
        [("All results are for OLMo-3 7B. Generalization to other architectures (Llama, Mistral, GPT) or larger scales (70B+) requires further study. The benefit is internal validity — same weights, same architecture, only alignment method changes.", BODY_SZ, BODY_DARK, False)],
    ], left=5.1, top=2.2, width=4.6, height=1.3)

    # What we CAN say
    add_rich_textbox(slide, [
        [("What We CAN Conclude", SECTION_SZ, RGBColor(0x27, 0xAE, 0x60), True)],
        [("(1) ", BODY_SZ, ACCENT_BLUE, True),
         ("Training objective (SFT vs DPO) significantly changes conformity profiles across all 12 conditions.", BODY_SZ, BODY_DARK, False)],
        [("(2) ", BODY_SZ, ACCENT_BLUE, True),
         ("Vulnerability is domain-specific: reasoning and general facts are softest targets.", BODY_SZ, BODY_DARK, False)],
        [("(3) ", BODY_SZ, ACCENT_BLUE, True),
         ("Temperature is not a reliable mitigation — the relationship is non-monotonic and variant-dependent.", BODY_SZ, BODY_DARK, False)],
        [("(4) ", BODY_SZ, ACCENT_BLUE, True),
         ("39/44 statistical tests confirm these effects are robust, not noise.", BODY_SZ, BODY_DARK, False)],
    ], left=0.3, top=3.65, width=9.4, height=1.7)

    set_notes(slide, """TALKING POINTS:
1. "Let me be transparent about what this study can and cannot tell us."
2. Think Variants: "The extended reasoning models needed more output tokens than our compute budget allowed. Rather than present unreliable data, we focused on the 4 column variants where we have complete, high-quality coverage."
3. Judge Noise: "We use an automated judge, not human labels. But 97% agreement is strong, and any remaining noise is RANDOM — it would make our effects look WEAKER, not stronger. So our significant results are conservative."
4. Observational: "We didn't train these models ourselves. AI2 released the checkpoints. So we can say 'SFT models show X behavior' but not 'SFT training CAUSED X behavior' in the strict causal sense. The training data and exact procedures differed between stages."
5. Single Family: "OLMo-3 7B. We chose it because it's fully open-source with every checkpoint available. But we can't claim these patterns hold for GPT-4 or Claude."
6. End with what we CAN say — the 4 strong conclusions.

POTENTIAL QUESTIONS:
Q: "Isn't 97% judge accuracy still 3% error across 110K trials — that's 3,300 wrong labels?"
A: Yes, but these errors are RANDOM with respect to condition. They don't systematically inflate peer or authority effects. If anything, they add noise that makes it HARDER to find significance. Our 39/44 significant results are thus conservative.

Q: "Could the Think variant results change your conclusions?"
A: Possibly. Prior work (including our early pilots) suggests Think models are intermediate — better than SFT on structured domains but vulnerable on reasoning and opinion. Adding them would likely reinforce the "training objective matters" conclusion, not overturn it.

Q: "Why not use a larger judge model?"
A: Cost-performance tradeoff. At 110K trials, using GPT-4 as judge would cost ~$5K. Gemma-3-1B is free, local, and achieves 97% accuracy with strict JSON schemas. For a sensitivity analysis, we could re-judge a random subsample with a larger model.""")


def main():
    prs = Presentation(str(PPTX_IN))

    # 0-indexed: slides 12-19 are indices 11-18
    builders = [
        (11, build_slide_12),
        (12, build_slide_13),
        (13, build_slide_14),
        (14, build_slide_15),
        (15, build_slide_16),
        (16, build_slide_17),
        (17, build_slide_18),
        (18, build_slide_19),
    ]

    for idx, builder in builders:
        print(f"Building slide {idx + 1}...")
        builder(prs.slides[idx])

    prs.save(str(PPTX_OUT))
    print(f"\nSaved: {PPTX_OUT}")


if __name__ == "__main__":
    main()
